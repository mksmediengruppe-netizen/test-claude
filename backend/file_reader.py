"""
File Reader v6.0 — Universal File Reader & Analyzer.

Reads and extracts content from any file type:
- Documents: PDF, DOCX, PPTX, TXT, MD, RTF
- Data: XLSX, CSV, TSV, JSON, XML
- Images: PNG, JPG, WEBP, GIF (with OCR via pytesseract)
- Archives: ZIP, TAR, GZ (recursive content listing)
- Code: .py, .js, .ts, .html, .css, .sql, .sh, .yaml, .toml, .ini
- Audio/Video: metadata extraction

Graceful degradation: if a library is missing, falls back to basic text extraction.
"""

import os
import io
import json
import csv
import zipfile
import tarfile
import logging
import mimetypes
import hashlib
from datetime import datetime
from pathlib import Path

logger = logging.getLogger("file_reader")

# Maximum file size: 50 MB
MAX_FILE_SIZE = 50 * 1024 * 1024
# Maximum text output length
MAX_TEXT_LENGTH = 100_000

# ══════════════════════════════════════════════════════════════
# Optional imports with graceful fallback
# ══════════════════════════════════════════════════════════════

_HAS_PYMUPDF = False
_HAS_DOCX = False
_HAS_PPTX = False
_HAS_OPENPYXL = False
_HAS_PANDAS = False
_HAS_PIL = False
_HAS_TESSERACT = False

try:
    import fitz  # PyMuPDF
    _HAS_PYMUPDF = True
except ImportError:
    pass

try:
    import docx
    _HAS_DOCX = True
except ImportError:
    pass

try:
    from pptx import Presentation
    _HAS_PPTX = True
except ImportError:
    pass

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    pass

try:
    import pandas as pd
    _HAS_PANDAS = True
except ImportError:
    pass

try:
    from PIL import Image
    _HAS_PIL = True
except ImportError:
    pass

try:
    import pytesseract
    _HAS_TESSERACT = True
except ImportError:
    pass


# Supported file extensions
SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.pptx', '.xlsx', '.xls',
    '.csv', '.tsv', '.json', '.xml',
    '.png', '.jpg', '.jpeg', '.gif', '.webp', '.svg',
    '.zip', '.tar', '.gz', '.tgz',
    '.txt', '.md', '.rst', '.rtf',
    '.py', '.js', '.ts', '.html', '.css', '.sql', '.sh', '.yaml', '.toml', '.ini',
    '.mp3', '.wav', '.mp4', '.webm', '.ogg',
}


class FileReadResult:
    """Result of reading a file."""

    def __init__(self, filename: str, file_type: str, size: int):
        self.filename = filename
        self.file_type = file_type
        self.size = size
        self.text = ""
        self.tables = []
        self.metadata = {}
        self.images_count = 0
        self.pages_count = 0
        self.error = None
        self.children = []  # For archives

    @property
    def success(self):
        return self.error is None

    @property
    def content(self):
        return self.text

    def to_dict(self):
        result = {
            "filename": self.filename,
            "file_type": self.file_type,
            "size": self.size,
            "size_human": _human_size(self.size),
            "text_length": len(self.text),
            "tables_count": len(self.tables),
            "images_count": self.images_count,
            "pages_count": self.pages_count,
            "metadata": self.metadata,
        }
        if self.error:
            result["error"] = self.error
        if self.children:
            result["archive_contents"] = len(self.children)
        return result

    def to_text(self, max_length=MAX_TEXT_LENGTH):
        """Format as text for LLM consumption."""
        parts = []
        parts.append(f"📄 File: {self.filename}")
        parts.append(f"   Type: {self.file_type} | Size: {_human_size(self.size)}")

        if self.pages_count:
            parts.append(f"   Pages: {self.pages_count}")
        if self.images_count:
            parts.append(f"   Images: {self.images_count}")
        if self.metadata:
            meta_str = ", ".join(f"{k}: {v}" for k, v in self.metadata.items() if v)
            if meta_str:
                parts.append(f"   Metadata: {meta_str}")

        if self.error:
            parts.append(f"\n⚠️ Error: {self.error}")

        if self.text:
            text = self.text[:max_length]
            if len(self.text) > max_length:
                text += f"\n\n... [truncated, {len(self.text)} chars total]"
            parts.append(f"\n--- Content ---\n{text}")

        if self.tables:
            parts.append(f"\n--- Tables ({len(self.tables)}) ---")
            for i, table in enumerate(self.tables[:5]):  # Max 5 tables
                parts.append(f"\nTable {i+1}:")
                parts.append(_format_table(table))

        if self.children:
            parts.append(f"\n--- Archive Contents ({len(self.children)} files) ---")
            for child in self.children[:20]:  # Max 20 entries
                parts.append(f"  • {child['name']} ({_human_size(child.get('size', 0))})")

        return "\n".join(parts)


def _human_size(size):
    """Convert bytes to human-readable size."""
    for unit in ["B", "KB", "MB", "GB"]:
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024
    return f"{size:.1f} TB"


def _format_table(table):
    """Format a table as text."""
    if not table or not table[0]:
        return "(empty table)"
    # Use first row as header
    header = table[0]
    lines = [" | ".join(str(c) for c in header)]
    lines.append("-" * len(lines[0]))
    for row in table[1:10]:  # Max 10 rows
        lines.append(" | ".join(str(c) for c in row))
    if len(table) > 11:
        lines.append(f"... ({len(table) - 1} rows total)")
    return "\n".join(lines)


def _detect_type(filepath):
    """Detect file type from extension and mime."""
    ext = Path(filepath).suffix.lower()
    mime, _ = mimetypes.guess_type(filepath)

    type_map = {
        # Documents
        ".pdf": "pdf",
        ".docx": "docx", ".doc": "doc",
        ".pptx": "pptx", ".ppt": "ppt",
        ".txt": "text", ".md": "markdown", ".rst": "text",
        ".rtf": "rtf",
        # Data
        ".xlsx": "xlsx", ".xls": "xls",
        ".csv": "csv", ".tsv": "tsv",
        ".json": "json",
        ".xml": "xml",
        # Images
        ".png": "image", ".jpg": "image", ".jpeg": "image",
        ".gif": "image", ".webp": "image", ".bmp": "image",
        ".svg": "svg",
        # Code
        ".py": "code", ".js": "code", ".ts": "code",
        ".html": "code", ".css": "code", ".scss": "code",
        ".sql": "code", ".sh": "code", ".bash": "code",
        ".yaml": "code", ".yml": "code", ".toml": "code",
        ".ini": "code", ".cfg": "code", ".conf": "code",
        ".java": "code", ".cpp": "code", ".c": "code",
        ".go": "code", ".rs": "code", ".rb": "code",
        ".php": "code", ".swift": "code", ".kt": "code",
        ".r": "code", ".m": "code",
        ".env": "code", ".gitignore": "code",
        ".dockerfile": "code",
        # Archives
        ".zip": "zip", ".tar": "tar",
        ".gz": "gzip", ".tgz": "tar",
        ".rar": "archive", ".7z": "archive",
        # Audio/Video
        ".mp3": "audio", ".wav": "audio", ".ogg": "audio",
        ".mp4": "video", ".avi": "video", ".mkv": "video",
        ".webm": "video",
    }

    return type_map.get(ext, "unknown")


# ══════════════════════════════════════════════════════════════
# Reader functions for each file type
# ══════════════════════════════════════════════════════════════

def _read_pdf(filepath, result):
    """Read PDF using PyMuPDF."""
    if not _HAS_PYMUPDF:
        result.error = "PyMuPDF not installed. Install with: pip install PyMuPDF"
        # Fallback: try to extract text with pdftotext
        try:
            import subprocess
            text = subprocess.run(
                ["pdftotext", filepath, "-"],
                capture_output=True, text=True, timeout=30
            ).stdout
            result.text = text
        except Exception:
            pass
        return

    doc = fitz.open(filepath)
    result.pages_count = len(doc)
    result.metadata = {
        "title": doc.metadata.get("title", ""),
        "author": doc.metadata.get("author", ""),
        "subject": doc.metadata.get("subject", ""),
        "creator": doc.metadata.get("creator", ""),
    }

    text_parts = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            text_parts.append(f"[Page {page_num + 1}]\n{text}")
        # Count images
        result.images_count += len(page.get_images())

    result.text = "\n\n".join(text_parts)

    # Extract tables if possible
    for page in doc:
        try:
            tables = page.find_tables()
            for table in tables:
                data = table.extract()
                if data:
                    result.tables.append(data)
        except Exception:
            pass

    doc.close()


def _read_docx(filepath, result):
    """Read DOCX using python-docx."""
    if not _HAS_DOCX:
        result.error = "python-docx not installed. Install with: pip install python-docx"
        return

    doc = docx.Document(filepath)
    result.metadata = {
        "author": doc.core_properties.author or "",
        "title": doc.core_properties.title or "",
        "created": str(doc.core_properties.created or ""),
    }

    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    result.text = "\n\n".join(text_parts)

    # Extract tables
    for table in doc.tables:
        data = []
        for row in table.rows:
            data.append([cell.text for cell in row.cells])
        if data:
            result.tables.append(data)

    result.images_count = len(doc.inline_shapes)


def _read_pptx(filepath, result):
    """Read PPTX using python-pptx."""
    if not _HAS_PPTX:
        result.error = "python-pptx not installed. Install with: pip install python-pptx"
        return

    prs = Presentation(filepath)
    result.pages_count = len(prs.slides)

    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text for cell in row.cells])
                result.tables.append(table_data)

        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

    result.text = "\n\n".join(text_parts)


def _read_xlsx(filepath, result):
    """Read XLSX using openpyxl and/or pandas."""
    if _HAS_PANDAS:
        try:
            xls = pd.ExcelFile(filepath)
            text_parts = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                text_parts.append(f"[Sheet: {sheet_name}] ({len(df)} rows x {len(df.columns)} cols)")
                text_parts.append(df.to_string(max_rows=50, max_cols=20))

                # Convert to table format
                header = list(df.columns)
                rows = df.head(50).values.tolist()
                result.tables.append([header] + rows)

            result.text = "\n\n".join(text_parts)
            result.pages_count = len(xls.sheet_names)
            return
        except Exception as e:
            logger.warning(f"pandas failed for xlsx: {e}")

    if _HAS_OPENPYXL:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        result.pages_count = len(wb.sheetnames)
        text_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=100, values_only=True):
                rows.append([str(c) if c is not None else "" for c in row])
            if rows:
                text_parts.append(f"[Sheet: {sheet_name}] ({ws.max_row} rows)")
                result.tables.append(rows)
                for row in rows[:50]:
                    text_parts.append(" | ".join(row))

        result.text = "\n\n".join(text_parts)
        wb.close()
    else:
        result.error = "openpyxl/pandas not installed"


def _read_csv(filepath, result, delimiter=","):
    """Read CSV/TSV files."""
    if _HAS_PANDAS:
        try:
            df = pd.read_csv(filepath, delimiter=delimiter, nrows=500)
            result.text = f"CSV: {len(df)} rows x {len(df.columns)} cols\n\n"
            result.text += df.to_string(max_rows=50, max_cols=20)
            header = list(df.columns)
            rows = df.head(50).values.tolist()
            result.tables.append([header] + rows)
            return
        except Exception:
            pass

    # Fallback: raw CSV reading
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = []
        for i, row in enumerate(reader):
            if i > 500:
                break
            rows.append(row)
        if rows:
            result.tables.append(rows)
            result.text = "\n".join(" | ".join(r) for r in rows[:50])


def _read_json(filepath, result):
    """Read JSON files."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        content = f.read(MAX_TEXT_LENGTH)

    try:
        data = json.loads(content)
        result.text = json.dumps(data, indent=2, ensure_ascii=False)[:MAX_TEXT_LENGTH]

        # If it's a list of dicts, extract as table
        if isinstance(data, list) and data and isinstance(data[0], dict):
            headers = list(data[0].keys())
            rows = [headers]
            for item in data[:50]:
                rows.append([str(item.get(h, "")) for h in headers])
            result.tables.append(rows)
    except json.JSONDecodeError:
        result.text = content


def _read_image(filepath, result):
    """Read image — extract metadata and OCR text."""
    result.images_count = 1

    if _HAS_PIL:
        try:
            img = Image.open(filepath)
            result.metadata = {
                "format": img.format,
                "size": f"{img.width}x{img.height}",
                "mode": img.mode,
            }
            # EXIF data
            if hasattr(img, "_getexif") and img._getexif():
                exif = img._getexif()
                if exif:
                    result.metadata["has_exif"] = True
        except Exception as e:
            logger.warning(f"PIL failed: {e}")

    if _HAS_TESSERACT and _HAS_PIL:
        try:
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img, lang="eng+rus")
            if text.strip():
                result.text = f"[OCR extracted text]\n{text.strip()}"
            else:
                result.text = "[No text detected in image via OCR]"
        except Exception as e:
            result.text = f"[OCR failed: {e}]"
    else:
        result.text = "[Image file. OCR not available — install pytesseract and Pillow for text extraction]"


def _read_text(filepath, result):
    """Read plain text files."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        result.text = f.read(MAX_TEXT_LENGTH)


def _read_code(filepath, result):
    """Read code files with syntax info."""
    ext = Path(filepath).suffix.lower()
    lang_map = {
        ".py": "python", ".js": "javascript", ".ts": "typescript",
        ".html": "html", ".css": "css", ".sql": "sql",
        ".sh": "bash", ".yaml": "yaml", ".yml": "yaml",
        ".json": "json", ".xml": "xml", ".java": "java",
        ".cpp": "c++", ".c": "c", ".go": "go", ".rs": "rust",
        ".rb": "ruby", ".php": "php", ".swift": "swift",
    }
    lang = lang_map.get(ext, "text")

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        content = f.read(MAX_TEXT_LENGTH)

    lines = content.count("\n") + 1
    result.metadata = {"language": lang, "lines": lines}
    result.text = f"```{lang}\n{content}\n```"


def _read_zip(filepath, result):
    """Read ZIP archive — list contents and extract text files."""
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            for info in zf.infolist():
                child = {
                    "name": info.filename,
                    "size": info.file_size,
                    "compressed": info.compress_size,
                    "is_dir": info.is_dir(),
                }
                result.children.append(child)

            # Try to read small text files
            text_parts = []
            for info in zf.infolist()[:10]:  # Max 10 files
                if info.is_dir() or info.file_size > 100_000:
                    continue
                ext = Path(info.filename).suffix.lower()
                if ext in (".txt", ".md", ".csv", ".json", ".py", ".js", ".html", ".css", ".yaml", ".yml"):
                    try:
                        content = zf.read(info.filename).decode("utf-8", errors="replace")
                        text_parts.append(f"--- {info.filename} ---\n{content[:5000]}")
                    except Exception:
                        pass

            if text_parts:
                result.text = "\n\n".join(text_parts)
    except zipfile.BadZipFile:
        result.error = "Invalid or corrupted ZIP file"


def _read_tar(filepath, result):
    """Read TAR/GZ archive."""
    try:
        mode = "r:gz" if filepath.endswith((".gz", ".tgz")) else "r"
        with tarfile.open(filepath, mode) as tf:
            for member in tf.getmembers():
                child = {
                    "name": member.name,
                    "size": member.size,
                    "is_dir": member.isdir(),
                }
                result.children.append(child)
    except Exception as e:
        result.error = f"Failed to read archive: {e}"


def _read_xml(filepath, result):
    """Read XML files."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        result.text = f.read(MAX_TEXT_LENGTH)


# ══════════════════════════════════════════════════════════════
# Main API
# ══════════════════════════════════════════════════════════════

def read_file(filepath: str) -> FileReadResult:
    """
    Read any file and extract its content.

    Args:
        filepath: Absolute path to the file.

    Returns:
        FileReadResult with extracted text, tables, metadata.
    """
    filename = os.path.basename(filepath)
    file_type = _detect_type(filepath)

    # Check file exists
    if not os.path.exists(filepath):
        result = FileReadResult(filename, "unknown", 0)
        result.error = f"File not found: {filepath}"
        return result

    size = os.path.getsize(filepath)

    # Check file size
    if size > MAX_FILE_SIZE:
        result = FileReadResult(filename, file_type, size)
        result.error = f"File too large ({_human_size(size)}). Maximum: {_human_size(MAX_FILE_SIZE)}"
        return result

    result = FileReadResult(filename, file_type, size)

    try:
        readers = {
            "pdf": _read_pdf,
            "docx": _read_docx,
            "doc": _read_docx,  # Will fail gracefully
            "pptx": _read_pptx,
            "xlsx": _read_xlsx,
            "xls": _read_xlsx,
            "csv": lambda fp, r: _read_csv(fp, r, ","),
            "tsv": lambda fp, r: _read_csv(fp, r, "\t"),
            "json": _read_json,
            "xml": _read_xml,
            "image": _read_image,
            "svg": _read_text,
            "text": _read_text,
            "markdown": _read_text,
            "code": _read_code,
            "zip": _read_zip,
            "tar": _read_tar,
            "gzip": _read_tar,
        }

        reader = readers.get(file_type)
        if reader:
            reader(filepath, result)
        else:
            # Fallback: try to read as text
            try:
                _read_text(filepath, result)
                result.file_type = "text (auto-detected)"
            except Exception:
                result.error = f"Unsupported file type: {file_type}"

    except Exception as e:
        result.error = f"Failed to read file: {str(e)}"
        logger.error(f"Error reading {filepath}: {e}", exc_info=True)

    return result


def get_supported_formats():
    """Return list of supported file formats with availability status."""
    return {
        "documents": {
            "pdf": {"available": _HAS_PYMUPDF, "library": "PyMuPDF"},
            "docx": {"available": _HAS_DOCX, "library": "python-docx"},
            "pptx": {"available": _HAS_PPTX, "library": "python-pptx"},
            "txt/md": {"available": True, "library": "built-in"},
        },
        "data": {
            "xlsx": {"available": _HAS_OPENPYXL or _HAS_PANDAS, "library": "openpyxl/pandas"},
            "csv/tsv": {"available": True, "library": "built-in"},
            "json": {"available": True, "library": "built-in"},
            "xml": {"available": True, "library": "built-in"},
        },
        "images": {
            "png/jpg/webp": {"available": _HAS_PIL, "library": "Pillow"},
            "ocr": {"available": _HAS_TESSERACT, "library": "pytesseract"},
        },
        "archives": {
            "zip": {"available": True, "library": "built-in"},
            "tar/gz": {"available": True, "library": "built-in"},
        },
        "code": {
            "all languages": {"available": True, "library": "built-in"},
        },
    }
